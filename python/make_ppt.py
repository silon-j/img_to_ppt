import locale
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# 使用系统 locale，匹配 Finder 的排序行为（中文按拼音排序）
try:
    locale.setlocale(locale.LC_ALL, '')
except locale.Error:
    pass

SOURCE_DIR = Path(__file__).parent.parent / "source"
OUTPUT_FILE = Path(__file__).parent.parent / "output.pptx"

# 幻灯片尺寸：16:9 宽屏
SLIDE_WIDTH = Emu(12192000)   # 10 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def finder_sort_key(name: str):
    """模拟 macOS Finder 排序：数字按大小、中文按拼音"""
    parts = re.split(r'(\d+)', name)
    return [int(p) if p.isdigit() else locale.strxfrm(p.lower()) for p in parts]


def collect_images(directory: Path) -> list[Path]:
    exts = {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff'}
    images = [f for f in directory.iterdir() if f.suffix.lower() in exts]
    images.sort(key=lambda p: finder_sort_key(p.name))
    return images


def add_image_slide(prs: Presentation, img_path: Path):
    """添加一页幻灯片，将图片居中并尽可能铺满页面（保持比例）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    with Image.open(img_path) as img:
        img_w, img_h = img.size

    slide_w = SLIDE_WIDTH
    slide_h = SLIDE_HEIGHT

    # 计算缩放比例，使图片在保持比例的前提下尽量铺满幻灯片
    ratio_w = slide_w / img_w
    ratio_h = slide_h / img_h
    ratio = min(ratio_w, ratio_h)

    pic_w = int(img_w * ratio)
    pic_h = int(img_h * ratio)

    # 居中
    left = (slide_w - pic_w) // 2
    top = (slide_h - pic_h) // 2

    slide.shapes.add_picture(str(img_path), left, top, pic_w, pic_h)


def main():
    images = collect_images(SOURCE_DIR)
    if not images:
        print("source 文件夹中没有找到图片。")
        return

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for img_path in images:
        print(f"添加: {img_path.name}")
        add_image_slide(prs, img_path)

    prs.save(str(OUTPUT_FILE))
    print(f"\n完成！共 {len(images)} 页，已保存到: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
